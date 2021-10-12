from pptx import Presentation
# from spellcheck import SpellChecker

prs = Presentation('ppts/test.pptx')

# text_runs will be populated with a list of strings,
# one for each text run in presentation
text_runs = {}
slide_count = 0

for slide in prs.slides:
    slide_count += 1
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            text_rendered = []
            for run in paragraph.runs:
                text_rendered.append(run)
    text_runs[slide_count] = text_rendered

for value in text_runs:
    print(key, value)


